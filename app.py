from flask import Flask, request, send_file
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from flask_cors import CORS
import io

app = Flask(__name__)
CORS(app)

@app.route('/generate', methods=['POST'])
def generate_ppt():
    print("✅ 收到請求了！")
    data = request.get_json()
    print("📄 歌詞內容：", data)

    title = data.get("title", "").strip()
    lyricist = data.get("lyricist", "").strip()
    composer = data.get("composer", "").strip()
    singer = data.get("singer", "").strip()
    lyrics = data.get("lyrics", "").strip()

    prs = Presentation()
    prs.slide_width = Inches(13.33)  # 16:9
    prs.slide_height = Inches(7.5)

    def add_cover_slide(title, lyricist, composer, singer):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(0, 0, 0)

        txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(3.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        content = []
        if title:
            content.append(title)
        if lyricist:
            content.append(f"作詞：{lyricist}")
        if composer:
            content.append(f"作曲：{composer}")
        if singer:
            content.append(f"原唱：{singer}")

        for line in content:
            p = tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = line
            run.font.size = Pt(60)
            run.font.bold = True
            run.font.name = 'Microsoft JhengHei'
            run.font.color.rgb = RGBColor(255, 255, 255)

    def add_lyrics_slide(text_lines):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(0, 0, 0)

        txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.33), Inches(3.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        for line in text_lines:
            p = tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = line
            run.font.size = Pt(60)
            run.font.bold = True
            run.font.name = 'Microsoft JhengHei'
            run.font.color.rgb = RGBColor(255, 255, 255)

    # 插入封面
    add_cover_slide(title, lyricist, composer, singer)

    # 分段：空行代表换页，每页最多 4 行
    blocks = []
    current_block = []

    for line in lyrics.splitlines():
        line = line.strip()
        if line == '':
            if current_block:
                blocks.append(current_block)
                current_block = []
        else:
            current_block.append(line)
            if len(current_block) == 4:
                blocks.append(current_block)
                current_block = []

    if current_block:
        blocks.append(current_block)

    for block in blocks:
        add_lyrics_slide(block)

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)

    return send_file(
        output,
        as_attachment=True,
        download_name=f"{title or 'lyrics'}.pptx",
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )

if __name__ == '__main__':
    app.run(debug=True)
